# SPDX-License-Identifier: Apache-2.0
"""llcore 研究知見スライド生成 (16:9 ワイド・シンプル・1 スライド 1 メッセージ + 説明図)。

ja / en の 2 本を python-pptx で生成。数値はすべて公開済み一次資料
(research/internalization_poc/ + research/paper/PAPER_DRAFT.md) の確定値。
図はネイティブ図形/チャートで描画 (画像依存なし・実データのみ、模式は明記)。
ライセンス: スライド本体 = CC BY 4.0 (出典明示で商用利用可)。

実行::  py -3.11 slides/make_slides.py
出力::  slides/llcore_findings_2026_ja.pptx / llcore_findings_2026_en.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

_HERE = Path(__file__).resolve().parent
ACCENT = RGBColor(0x4F, 0x81, 0xBD)
ACCENT_LT = RGBColor(0xDC, 0xE6, 0xF1)
RED_LT = RGBColor(0xD9, 0x96, 0x94)
GREEN = RGBColor(0x9B, 0xBB, 0x59)
PURPLE = RGBColor(0x80, 0x64, 0xA2)
ORANGE = RGBColor(0xE8, 0xA3, 0x3D)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
GRAY_BAR = RGBColor(0x7F, 0x7F, 0x7F)
W, H = Inches(13.333), Inches(7.5)   # 16:9

FOOTER = {
    "ja": "© 2026 Kazufumi Furuse — CC BY 4.0 | 出典: github.com/furuse-kazufumi/llcore | qiita.com/furuse-kazufumi",
    "en": "© 2026 Kazufumi Furuse — CC BY 4.0 | Source: github.com/furuse-kazufumi/llcore | qiita.com/furuse-kazufumi",
}

# 本文 (左カラム) と図 (右カラム) の標準レイアウト
TXT_X, TXT_W = Inches(0.9), Inches(6.1)
FIG_X, FIG_W = Inches(7.3), Inches(5.5)
BODY_Y, BODY_H = Inches(1.7), Inches(4.7)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _txt(slide, x, y, w, h, lines, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_after=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines if isinstance(lines, list) else [lines]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            text, lv = line
            p.level = lv
        else:
            text = line
        p.text = text
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Meiryo UI"
    return box


def _box(slide, x, y, w, h, text, fill=ACCENT_LT, font=11, bold=False, color=DARK,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(font)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Meiryo UI"
    return sp


def _arrow(slide, x1, y1, x2, y2, color=GRAY, weight=2.25):
    # Length / 2 等の演算は float 化し PowerPoint が XML を拒否する → 必ず int (EMU) に丸める
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(weight)
    # 矢じり (a:tailEnd) — lxml SubElement で親 nsmap を継承させる (makeelement は PPT が拒否)
    from lxml import etree
    from pptx.oxml.ns import qn
    ln = cn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "arrow")
    return cn


def _bar_chart(slide, x, y, w, h, categories, values, colors, axis_title="",
               number_format='0.0', font_pt=11, min_scale=None, max_scale=None):
    """単系列の縦棒チャート (実データのみ)。colors は棒ごと。"""
    data = CategoryChartData()
    data.categories = categories
    data.add_series("v", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    ch.font.size = Pt(font_pt)
    ch.font.name = "Meiryo UI"
    plot = ch.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.font.size = Pt(font_pt)
    dl.font.bold = True
    for i, pt_ in enumerate(plot.series[0].points):
        pt_.format.fill.solid()
        pt_.format.fill.fore_color.rgb = colors[i % len(colors)]
    va = ch.value_axis
    va.has_major_gridlines = False
    va.visible = False
    if min_scale is not None:
        va.minimum_scale = min_scale
    if max_scale is not None:
        va.maximum_scale = max_scale
    ch.category_axis.format.line.color.rgb = GRAY
    return gf


def _fig_caption(slide, text, x=FIG_X, w=FIG_W, y=Inches(6.0)):
    _txt(slide, x, y, w, Inches(0.5), text, size=10.5, color=GRAY)


def _slide(prs, lang, title, bullets, note=None, body_w=TXT_W):
    s = _blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.62), Inches(0.12),
                             Inches(0.62))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _txt(s, Inches(0.85), Inches(0.5), Inches(11.9), Inches(1.0), title, size=29, bold=True)
    _txt(s, TXT_X, BODY_Y, body_w, BODY_H, bullets, size=17, space_after=12)
    if note:
        _txt(s, Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.45), note, size=11, color=GRAY)
    _txt(s, Inches(0.9), Inches(7.0), Inches(11.6), Inches(0.4), FOOTER[lang], size=10, color=GRAY)
    return s


# ---------------------------------------------------------------------------
# 説明図 (スライド別)
# ---------------------------------------------------------------------------

def fig_gate_compare(s, ja):
    """S2: 経験的ゲート vs 証明ゲートの位置の違い (フロー対比図)。"""
    y1, y2 = Inches(2.0), Inches(3.9)
    bw, bh, gap = Inches(1.45), Inches(0.75), Inches(0.30)
    x0 = FIG_X
    _txt(s, x0, y1 - Inches(0.42), FIG_W, Inches(0.35),
         "経験的ゲート (DGM/SEAL 系)" if ja else "Empirical gate (DGM/SEAL line)",
         size=13, bold=True, color=GRAY)
    _box(s, x0, y1, bw, bh, "変異" if ja else "mutate")
    _arrow(s, x0 + bw, y1 + bh / 2, x0 + bw + gap, y1 + bh / 2)
    _box(s, x0 + bw + gap, y1, bw, bh, "評価 = 実行\n(死を被る)" if ja else "evaluate = run\n(pays deaths)",
         fill=RED_LT)
    _arrow(s, x0 + 2 * (bw + gap) - gap, y1 + bh / 2, x0 + 2 * (bw + gap), y1 + bh / 2)
    _box(s, x0 + 2 * (bw + gap), y1, bw, bh,
         "事後スコア判定\n(Goodhart 可)" if ja else "post-hoc score\n(Goodhart-able)", fill=RED_LT)

    _txt(s, x0, y2 - Inches(0.42), FIG_W, Inches(0.35),
         "証明ゲート (llcore)" if ja else "Proof gate (llcore)",
         size=13, bold=True, color=ACCENT)
    _box(s, x0, y2, bw, bh, "変異" if ja else "mutate")
    _arrow(s, x0 + bw, y2 + bh / 2, x0 + bw + gap, y2 + bh / 2)
    _box(s, x0 + bw + gap, y2, bw, bh, "証明 ρ<1\n(評価の前)" if ja else "prove ρ<1\n(before eval)",
         fill=ACCENT, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    _arrow(s, x0 + 2 * (bw + gap) - gap, y2 + bh / 2, x0 + 2 * (bw + gap), y2 + bh / 2)
    _box(s, x0 + 2 * (bw + gap), y2, bw, bh, "評価\n(死なない)" if ja else "evaluate\n(no deaths)",
         fill=GREEN)
    _txt(s, x0, y2 + bh + Inches(0.15), FIG_W, Inches(0.4),
         "証明できない変更は評価前に却下 (fail-closed)" if ja else
         "Unprovable changes rejected before evaluation (fail-closed)", size=12, color=GRAY)


def fig_ladder(s, ja):
    """S3: certifier ladder 3 段。"""
    x0, w = FIG_X + Inches(0.3), Inches(4.3)
    labels = [
        ("閉形式 ∞-norm — O(n²)・ソルバ不要" if ja else "closed-form ∞-norm — O(n²), no solver", ACCENT_LT),
        ("頂点 2-norm (SVD)" if ja else "vertex 2-norm (SVD)", ACCENT_LT),
        ("Lyapunov SDP (CLARABEL)" if ja else "Lyapunov SDP (CLARABEL)", ACCENT_LT),
    ]
    for i, (txt, fill) in enumerate(labels):
        _box(s, x0 + Inches(0.25) * i, Inches(2.1) + Inches(0.95) * i, w, Inches(0.7), txt,
             fill=fill, font=13)
    _box(s, x0 + Inches(0.5), Inches(5.05), w, Inches(0.7),
         "証明なし → 不採用 (fail-closed)" if ja else "no proof → rejected (fail-closed)",
         fill=ACCENT, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, font=13)
    _txt(s, x0 - Inches(0.25), Inches(1.75), Inches(5), Inches(0.35),
         "健全 certifier の 3 段ラダー (下ほど強い)" if ja else
         "Three sound certifier rungs (stronger downward)", size=12, color=GRAY)


def fig_kappa_mechanisms(s, ja):
    """S4: κ 段差 + 3 機構の箱。"""
    # κ step (折れ線を connector 2 本で)
    x0, y_hi, y_lo = FIG_X, Inches(2.3), Inches(2.9)
    xm = x0 + Inches(2.4)
    x1 = x0 + Inches(5.0)
    for (a, b, c, d) in [(x0, y_lo, xm, y_lo), (xm, y_lo, xm, y_hi), (xm, y_hi, x1, y_hi)]:
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a, b, c, d)
        cn.line.color.rgb = ACCENT
        cn.line.width = Pt(3)
    _txt(s, x0, y_lo + Inches(0.05), Inches(2.2), Inches(0.3), "κ = 1.0 (安全)" if ja else "κ = 1.0 (safe)",
         size=12, color=GRAY)
    _txt(s, xm + Inches(0.1), y_hi - Inches(0.35), Inches(3.2), Inches(0.3),
         "κ = 2.0 — 旧安全個体が発散 = 死" if ja else "κ = 2.0 — old-safe genes diverge = death",
         size=12, color=DARK, bold=True)
    # 3 機構
    bw, bh = Inches(1.7), Inches(1.15)
    items = [
        ("ENDO\n" + ("自己予見 (証明)" if ja else "foresight (proof)"), ACCENT),
        ("REVIVE\n" + ("死→修復→復活" if ja else "die→repair→revive"), PURPLE),
        ("OBSERVE\n" + ("他者の死から学ぶ" if ja else "learn from others' deaths"), ORANGE),
    ]
    for i, (txt, fill) in enumerate(items):
        _box(s, FIG_X + (bw + Inches(0.18)) * i, Inches(3.9), bw, bh, txt, fill=fill,
             color=RGBColor(0xFF, 0xFF, 0xFF), font=12, bold=True)
    _txt(s, FIG_X, Inches(5.2), FIG_W, Inches(0.6),
         "+ baseline: NONE (無防備) / EXO (設計時固定ゲート)" if ja else
         "+ baselines: NONE (no gate) / EXO (design-time frozen gate)", size=12, color=GRAY)


def fig_deaths_bar(s, ja):
    """S5: 定常死の実データ棒グラフ。"""
    cats = ["NONE\n" + ("無防備" if ja else "no gate"),
            "OBSERVE\n" + ("観察学習" if ja else "observational"),
            "EXO\n" + ("固定ゲート" if ja else "frozen gate"),
            "ENDO\n" + ("証明" if ja else "proof")]
    _bar_chart(s, FIG_X, Inches(1.9), FIG_W, Inches(3.9), cats,
               (27.2, 17.3, 10.2, 0.0), [GRAY_BAR, ORANGE, RED_LT, ACCENT])
    _fig_caption(s, "致命的評価数 (定常状態, linear 基質, n=20 平均) — 一次データ: results_viability_ab.json"
                 if ja else
                 "Lethal evaluations (steady state, linear substrate, n=20 mean) — data: results_viability_ab.json")


def fig_memory_bar(s, ja):
    """S6: 記憶保存 (pop_mean) の実データ棒グラフ。"""
    cats = ["NONE\n" + ("死=消滅" if ja else "death = erasure"),
            "REVIVE\n" + ("死→復活" if ja else "die → revive"),
            "ENDO\n" + ("証明" if ja else "proof")]
    _bar_chart(s, FIG_X, Inches(1.9), FIG_W, Inches(3.9), cats,
               (0.649, 0.708, 0.750), [GRAY_BAR, PURPLE, ACCENT],
               number_format='0.000', min_scale=0.6, max_scale=0.78)
    _fig_caption(s, "集団平均適応度 (linear 基質)。縦軸は 0.60 開始 (差の可視化のため; 効果量に注意)"
                 if ja else
                 "Population-mean fitness (linear). Y-axis starts at 0.60 (to make the gap visible; mind the effect size)")


def fig_factorial_bar(s, ja):
    """S7: factorial 組合せの実データ棒グラフ。"""
    cats = ["NONE", "O", "R", "R+O", ("E 入り 4 組" if ja else "all 4 E-combos")]
    _bar_chart(s, FIG_X, Inches(1.9), FIG_W, Inches(3.9), cats,
               (27.2, 17.3, 18.8, 6.2, 0.0), [GRAY_BAR, ORANGE, PURPLE, GREEN, ACCENT])
    _fig_caption(s, "致命的評価数 (2³ factorial, linear, n=20 平均)。E=証明 / R=復活 / O=観察"
                 if ja else
                 "Lethal evaluations (2^3 factorial, linear, n=20 mean). E=proof / R=revival / O=observation")


def fig_trust_bar(s, ja):
    """S8: trust の regime 追従 (実データ 2 本) + premise 破壊の説明。"""
    _box(s, FIG_X, Inches(1.95), FIG_W, Inches(0.85),
         "検証器が κ を 40% 過小推定 (premise 破壊)\n→ 証明が嘘をつく regime を人工的に構成" if ja else
         "Verifier under-senses κ by 40% (premise broken)\n→ a regime where the proof lies",
         fill=RED_LT, font=12)
    cats = [("GOOD 区間\n(前提 正常)" if ja else "GOOD blocks\n(premise holds)"),
            ("BAD 区間\n(前提 破壊)" if ja else "BAD blocks\n(premise broken)")]
    _bar_chart(s, FIG_X + Inches(0.6), Inches(3.0), Inches(4.3), Inches(2.7), cats,
               (0.91, 0.51), [GREEN, RED_LT], number_format='0.00', max_scale=1.0)
    _fig_caption(s, "META の証明への信頼度 (平均, linear)。「許可したのに死んだ」だけから regime を追従"
                 if ja else
                 "META's trust in the proof (mean, linear). Tracks the regime from admitted-then-died evidence alone")


def fig_quadrant(s, ja):
    """S9: 2×2 マトリクス (保証の種類 × 対象)。"""
    x0, y0 = FIG_X + Inches(0.45), Inches(2.0)
    qw, qh = Inches(2.3), Inches(1.7)
    # 軸ラベル
    _txt(s, x0 - Inches(0.45), y0 - Inches(0.35), Inches(5.2), Inches(0.3),
         ("対象 ↑ 連続な内部力学 / ↓ 離散タスク・成果物" if ja else
          "object: ↑ continuous internal dynamics / ↓ discrete tasks & artifacts"),
         size=11, color=GRAY)
    _txt(s, x0 - Inches(0.45), y0 + 2 * qh + Inches(0.1), Inches(5.2), Inches(0.3),
         ("保証: ← 統計的/経験的 | 健全な証明 (演繹) →" if ja else
          "guarantee: ← statistical/empirical | sound deductive proof →"),
         size=11, color=GRAY)
    # 4 象限
    _box(s, x0, y0, qw, qh, ("(空白)\n統計 × 連続力学" if ja else "(open)\nstatistical x continuous"),
         fill=RGBColor(0xF2, 0xF2, 0xF2), font=11, color=GRAY)
    _box(s, x0 + qw + Inches(0.12), y0, qw, qh,
         "llcore (PoC)\n" + ("証明 × 連続記憶力学" if ja else "proof x continuous memory dynamics"),
         fill=ACCENT, color=RGBColor(0xFF, 0xFF, 0xFF), font=12, bold=True)
    _box(s, x0, y0 + qh + Inches(0.12), qw, qh,
         "SGM (統計的GM)\nPAC limits (Wang+)" if ja else "SGM (statistical GM)\nPAC limits (Wang+)",
         fill=RGBColor(0xF2, 0xF2, 0xF2), font=11)
    _box(s, x0 + qw + Inches(0.12), y0 + qh + Inches(0.12), qw, qh,
         "SEVerA\n" + ("Dafny 論理契約" if ja else "Dafny logical contracts"),
         fill=ACCENT_LT, font=11)
    _fig_caption(s, "3 論文とも本文を一次取得して確認 (敵対検証付き)" if ja else
                 "All three verified against primary sources (with adversarial review)", y=Inches(5.85))


# ---------------------------------------------------------------------------

def build(lang: str) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    ja = lang == "ja"

    # 1. タイトル
    s = _blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.6), W, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _txt(s, Inches(1.2), Inches(2.0), Inches(11), Inches(1.4),
         "自己進化 AI を「証明」で守る" if ja else "Guarding Self-Evolving AI with Proofs",
         size=44, bold=True)
    _txt(s, Inches(1.2), Inches(3.4), Inches(11), Inches(0.8),
         "llcore 研究知見 2026 — 検証付きニューラル進化 (CPU 完結)" if ja else
         "llcore research findings 2026 — verified neural evolution on CPU",
         size=24, color=ACCENT)
    _txt(s, Inches(1.2), Inches(4.9), Inches(11), Inches(1.2), [
        "Kazufumi Furuse",
        "github.com/furuse-kazufumi/llcore (paper draft + 全実験コード/データ公開)" if ja else
        "github.com/furuse-kazufumi/llcore (paper draft + all experiment code/data public)",
        "本スライド: CC BY 4.0 — 出典明示で商用利用可" if ja else
        "These slides: CC BY 4.0 — commercial use permitted with attribution",
    ], size=16, color=GRAY)
    _txt(s, Inches(1.2), Inches(7.0), Inches(11), Inches(0.4), FOOTER[lang], size=10, color=GRAY)

    # 2. 問題
    s = _slide(prs, lang,
               "問題: 自己進化 AI は何を根拠に自分の変更を許可するか" if ja else
               "Problem: on what grounds may a self-evolving AI admit its own changes?",
               [
                   "AI が自分のコア (記憶・力学) を書き換えながら進化する時代" if ja else
                   "AI now evolves its own core (memory, dynamics)",
                   "主流 (DGM / SEAL 系) = 経験的ゲート: スコアや過去の失敗で事後判定" if ja else
                   "Mainstream (DGM / SEAL) = empirical gates: post-hoc scores & past failures",
                   "構造的弱点: ① 死んで学ぶ ② Goodhart 可能 (指標ハック)" if ja else
                   "Structural weakness: (1) learns by dying (2) Goodhart-able",
                   "llcore: 性質を絞れば「数学的証明」が評価の前のゲートになる" if ja else
                   "llcore: narrow the property, and a proof can gate before evaluation",
               ])
    fig_gate_compare(s, ja)

    # 3. アプローチ
    s = _slide(prs, lang,
               "アプローチ: 健全な収縮証明ゲート (fail-closed)" if ja else
               "Approach: a sound contraction-certifier gate (fail-closed)",
               [
                   "証明する性質 = 収縮 (ρ < 1): 内部状態が発散しない" if ja else
                   "Certified property = contraction (ρ < 1): internal state cannot diverge",
                   "証明は定理の帰結 = 捏造不能 (スコアと違いハックできない)" if ja else
                   "Proofs are theorem-derived = non-fabricable (unlike scores)",
                   "進化ループ内・評価前に gate — 事後検証より 17-19 倍安い" if ja else
                   "Gated in-loop, before evaluation — 17-19x cheaper than post-hoc",
                   "全実験は事前登録 → 結果の順。negative も全強度で開示" if ja else
                   "All experiments pre-registered; negatives disclosed at full strength",
               ])
    fig_ladder(s, ja)

    # 4. 実験場
    s = _slide(prs, lang,
               "実験場: 死ねる環境 × 経験が記憶になる 3 つの道" if ja else
               "Testbed: a lethal environment x three roads from experience to memory",
               [
                   "環境が再帰ゲイン κ を 2 倍にステップ → 旧安全個体が発散 = 死" if ja else
                   "The environment steps κ by 2x → previously-safe genes diverge = death",
                   "死は回避不能な実在の脅威 (入力側の工夫では逃げられない設計)" if ja else
                   "Death is a real, unavoidable threat (cannot be routed around)",
                   "測定は助走後の定常状態 (n=20 シード, 事前登録)" if ja else
                   "Measured in the steady state after warm-up (n=20 seeds, pre-registered)",
               ])
    fig_kappa_mechanisms(s, ja)

    # 5. 知見 1
    s = _slide(prs, lang,
               "知見 1: 証明は死ぬ前に分かる — 定常死ゼロ" if ja else
               "Finding 1: proofs know before dying — zero steady-state deaths",
               [
                   "証明 (ENDO) は致命的評価を完全にゼロ化" if ja else
                   "The proof (ENDO) eliminates lethal evaluations entirely",
                   "観察学習は機能する — が、学んでも死が残る (17.3)。「死んで学ぶ」は構造コスト" if ja else
                   "Observational learning works — but deaths remain (17.3). Learning-by-dying is structural",
                   "12/12 の事前登録環境 (シード/κ/次元/課題難度) で頑健" if ja else
                   "Robust in 12/12 pre-registered configurations (seed/κ/dimension/difficulty)",
                   "記憶獲得能力は犠牲にならない (最大適応度差 0.7%)" if ja else
                   "Memory capability is not sacrificed (worst fitness gap 0.7%)",
               ],
               note="p < 0.001 (paired sign-flip, n=20)" if ja else "p < 0.001 (paired sign-flip, n=20)")
    fig_deaths_bar(s, ja)

    # 6. 知見 2
    s = _slide(prs, lang,
               "知見 2: 復活がないと、経験は記憶にならない" if ja else
               "Finding 2: without revival, experience does not become memory",
               [
                   "死 = 個体の除去 = 蓄積した経験 (遺伝子) も消える" if ja else
                   "Death removes the individual — and erases its accumulated experience",
                   "REVIVE は無防備と同程度死ぬのに、集団記憶が有意に高い" if ja else
                   "REVIVE dies about as often as no-gate, yet preserves more population memory",
                   "Δ = +0.060 (p = 0.0011) — 修復が記憶チャネルを保ったまま力学を安全化" if ja else
                   "Δ = +0.060 (p = 0.0011) — repair preserves the memory channel while safing the dynamics",
                   "死が「消滅」でなく「傷」になる" if ja else "Death becomes a scar, not a void",
               ])
    fig_memory_bar(s, ja)

    # 7. 知見 3
    s = _slide(prs, lang,
               "知見 3: 証明 1 つは、経験的機構の積み重ねに勝る" if ja else
               "Finding 3: one proof beats stacked empirical mechanisms",
               [
                   "3 機構の全 8 組合せを同条件で比較 (2³ factorial)" if ja else
                   "All 8 combinations compared under identical conditions (2^3 factorial)",
                   "証明 (E) を含む 4 組合せ: すべて死 0.0" if ja else
                   "All 4 proof-containing combos: zero deaths",
                   "証明なしの最良 = 復活+観察 (6.2)。単独より良いが、0 には届かない" if ja else
                   "Best proof-free combo = revival+observation (6.2) — better than alone, never zero",
                   "証明があると経験的機構は完全に冗長化 (交互作用 p=0.0003/0.0045)" if ja else
                   "With the proof, empirical mechanisms become redundant (interaction p=0.0003/0.0045)",
               ])
    fig_factorial_bar(s, ja)

    # 8. 知見 4
    s = _slide(prs, lang,
               "知見 4 (限界): 証明は前提が正しい時だけ健全" if ja else
               "Finding 4 (limits): a proof is only sound while its premises hold",
               [
                   "検証器の環境認識を壊すと、証明は致命域を「安全」と誤認する" if ja else
                   "Break the verifier's sensing, and the certificate admits a lethal band",
                   "観測可能な証拠だけで証明⇔経験を切替える信頼度学習は機能する (図)" if ja else
                   "A trust controller on observable evidence does track the regime (figure)",
                   "ただし得をするのは前提違反が高くつく環境のみ — 仮説の一部不成立をそのまま報告" if ja else
                   "But hedging pays only where violations are costly — a failed hypothesis is reported as-is",
                   "次の問い: 証明の前提を誰が守るのか (premise monitoring)" if ja else
                   "Open question: who certifies the certifier's premises?",
               ],
               note="PoC スケール (スカラー遺伝子, n=20) — 実 LLM スケールは未証明" if ja else
                    "PoC scale (scalar genes, n=20) — not yet shown at real-LLM scale")
    fig_trust_bar(s, ja)

    # 9. 研究地図
    s = _slide(prs, lang,
               "研究地図: 最近接研究との違い (一次ソース検証済み)" if ja else
               "Map: how this differs from the nearest work (primary sources verified)",
               [
                   "SGM: 統計的 certificate × 離散タスク性能 — 確率保証・捏造耐性なし" if ja else
                   "SGM: statistical certificate over discrete task scores — probabilistic",
                   "SEVerA: Dafny 論理契約の sound 検証 — 対象は離散 I/O" if ja else
                   "SEVerA: sound Dafny logical contracts — discrete I/O artifacts",
                   "PAC limits (Wang+): 汎化保証 — 安定性の保証ではない" if ja else
                   "PAC limits (Wang+): generalization bounds — not stability",
                   "右上の象限 (証明 × 連続記憶力学) を llcore が PoC 占有" if ja else
                   "llcore occupies the top-right quadrant (proof x continuous memory dynamics) at PoC scale",
               ])
    fig_quadrant(s, ja)

    # 10. 利用条件
    s = _slide(prs, lang, "利用について" if ja else "Using this work",
               [
                   "本スライド: CC BY 4.0 — 出典明示で、企業内研修・技術調査・製品検討に利用可" if ja else
                   "These slides: CC BY 4.0 — corporate training, tech scouting, product evaluation OK with attribution",
                   "コード/論文 draft: Apache-2.0 + Commercial dual license" if ja else
                   "Code / paper draft: Apache-2.0 + Commercial dual license",
                   "クローズド統合・SLA・NDA 相談 = 商用窓口: kazufumi@furuse.work" if ja else
                   "Closed-source integration / SLA / NDA consulting: kazufumi@furuse.work",
                   "出典表記例: 「Kazufumi Furuse, llcore (2026), github.com/furuse-kazufumi/llcore」" if ja else
                   "Attribution: \"Kazufumi Furuse, llcore (2026), github.com/furuse-kazufumi/llcore\"",
               ], body_w=Inches(6.6))
    # 右に出典 QR 的なリンクボックス
    _box(s, FIG_X + Inches(0.4), Inches(2.2), Inches(4.6), Inches(0.9),
         "github.com/furuse-kazufumi/llcore", fill=ACCENT, color=RGBColor(0xFF, 0xFF, 0xFF),
         font=15, bold=True)
    _box(s, FIG_X + Inches(0.4), Inches(3.3), Inches(4.6), Inches(0.9),
         "qiita.com/furuse-kazufumi", fill=ACCENT_LT, font=15)
    _box(s, FIG_X + Inches(0.4), Inches(4.4), Inches(4.6), Inches(0.9),
         ("商用ライセンス窓口\nkazufumi@furuse.work" if ja else
          "Commercial licensing\nkazufumi@furuse.work"), fill=ACCENT_LT, font=14)

    out = _HERE / f"llcore_findings_2026_{lang}.pptx"
    prs.save(str(out))
    return out


def main():
    for lang in ("ja", "en"):
        out = build(lang)
        print(f"wrote {out} ({out.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
